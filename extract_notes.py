"""Extract speaker notes from a PPTX file.

Usage:
    python extract_notes.py <pptx_path> <output_json>

Outputs a JSON array of {"slide": N, "text": "..."} objects.
Slides with no notes have empty text.
"""

from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation


ROBOTIC_PATTERNS = [
    "click on", "click the", 'click "',
    "navigate to", "select the", "go to",
    "enter the", "in the", "tap on", "tap the",
    "press the", "choose the", "open the",
]


def extract_notes(pptx_path: Path) -> list[dict]:
    """Extract speaker notes from each slide."""
    prs = Presentation(str(pptx_path))
    notes = []
    for i, slide in enumerate(prs.slides, 1):
        text = ""
        if slide.has_notes_slide:
            tf = slide.notes_slide.notes_text_frame
            text = tf.text.strip()
        notes.append({"slide": i, "text": text})
    return notes


def needs_refinement(notes: list[dict]) -> bool:
    """Auto-detect robotic documentation-style notes (Guidde, Scribe, Tango)."""
    non_empty = [n for n in notes if n["text"].strip()]
    if not non_empty:
        return False
    robotic = sum(
        1
        for n in non_empty
        if any(n["text"].strip().lower().startswith(p) for p in ROBOTIC_PATTERNS)
    )
    return robotic / len(non_empty) >= 0.6


def main():
    if len(sys.argv) < 3:
        print(f"Usage: python {sys.argv[0]} <pptx_path> <output_json>")
        sys.exit(1)

    pptx_path = Path(sys.argv[1])
    output_path = Path(sys.argv[2])

    if not pptx_path.exists():
        print(f"Error: PPTX not found: {pptx_path}")
        sys.exit(1)

    notes = extract_notes(pptx_path)

    # Summary
    total = len(notes)
    with_notes = sum(1 for n in notes if n["text"])
    without_notes = total - with_notes
    print(f"Total slides: {total}")
    print(f"  With notes: {with_notes}")
    print(f"  Without notes: {without_notes}")

    # Robotic detection
    if needs_refinement(notes):
        print("  Robotic style detected — refinement recommended")
    else:
        print("  Notes appear natural — refinement not needed")

    # Preview
    for n in notes:
        preview = n["text"][:80] + "..." if len(n["text"]) > 80 else n["text"]
        print(f"  Slide {n['slide']}: {preview if n['text'] else '(no notes)'}")

    # Save
    output_path.parent.mkdir(parents=True, exist_ok=True)
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(notes, f, indent=2, ensure_ascii=False)
    print(f"\nSaved to: {output_path}")


if __name__ == "__main__":
    main()
