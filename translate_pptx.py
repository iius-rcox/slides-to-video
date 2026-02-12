#!/usr/bin/env python3
"""Automated PPTX translation pipeline using Anthropic messages API.

Features:
- Deterministic traversal/write-back order
- Character-budget batching with retry handling
- Stable machine-readable translation report schema
- Glossary + never-translate enforcement
- Preflight checks for required artifacts
"""

from __future__ import annotations

import argparse
import json
import os
import re
import time
import urllib.error
import urllib.request
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

RUN_SEPARATOR = "||N||"
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_DIAGRAM_DATA_RELTYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramData"


@dataclass(frozen=True)
class TextItem:
    id: int
    text: str
    kind: str
    role: str
    slide_number: int
    location: str
    run_count: int


def _shape_role(shape: Any) -> str:
    if shape.is_placeholder:
        ph_type = shape.placeholder_format.type
        if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
            return "title"
        if ph_type == PP_PLACEHOLDER.SUBTITLE:
            return "subtitle"
    shape_name = (shape.name or "").lower()
    if "title" in shape_name:
        return "title"
    if "subtitle" in shape_name:
        return "subtitle"
    return "body"


def _nonempty_run_texts(paragraph: Any) -> list[str]:
    return [run.text for run in paragraph.runs if run.text and run.text.strip()]


def collect_text_items(prs: Presentation) -> list[TextItem]:
    items: list[TextItem] = []
    idx = 0

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            role = _shape_role(shape)
            if shape.has_text_frame:
                for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    run_texts = _nonempty_run_texts(paragraph)
                    if not run_texts:
                        continue
                    items.append(
                        TextItem(
                            id=idx,
                            text=RUN_SEPARATOR.join(run_texts),
                            kind="slide",
                            role=role,
                            slide_number=slide_idx,
                            location=f"slide={slide_idx};shape={shape.shape_id};para={para_idx}",
                            run_count=len(run_texts),
                        )
                    )
                    idx += 1

            if shape.has_table:
                for row_idx, row in enumerate(shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        for para_idx, paragraph in enumerate(cell.text_frame.paragraphs):
                            run_texts = _nonempty_run_texts(paragraph)
                            if not run_texts:
                                continue
                            items.append(
                                TextItem(
                                    id=idx,
                                    text=RUN_SEPARATOR.join(run_texts),
                                    kind="table",
                                    role="body",
                                    slide_number=slide_idx,
                                    location=f"slide={slide_idx};shape={shape.shape_id};table={row_idx},{col_idx};para={para_idx}",
                                    run_count=len(run_texts),
                                )
                            )
                            idx += 1

        rels = sorted(slide.part.rels.values(), key=lambda rel: rel.rId)
        for rel in rels:
            if rel.reltype != _DIAGRAM_DATA_RELTYPE:
                continue
            dgm_root = etree.fromstring(rel.target_part.blob)
            t_elements = list(dgm_root.iter(f"{{{_A_NS}}}t"))
            for t_idx, element in enumerate(t_elements):
                if element.text and element.text.strip():
                    items.append(
                        TextItem(
                            id=idx,
                            text=element.text,
                            kind="smartart",
                            role="body",
                            slide_number=slide_idx,
                            location=f"slide={slide_idx};smartart_rel={rel.rId};text={t_idx}",
                            run_count=1,
                        )
                    )
                    idx += 1

        if slide.has_notes_slide:
            for para_idx, paragraph in enumerate(slide.notes_slide.notes_text_frame.paragraphs):
                run_texts = _nonempty_run_texts(paragraph)
                if not run_texts:
                    continue
                items.append(
                    TextItem(
                        id=idx,
                        text=RUN_SEPARATOR.join(run_texts),
                        kind="notes",
                        role="body",
                        slide_number=slide_idx,
                        location=f"slide={slide_idx};notes_para={para_idx}",
                        run_count=len(run_texts),
                    )
                )
                idx += 1

    return items


def build_batches(items: list[TextItem], max_chars: int) -> list[list[TextItem]]:
    batches: list[list[TextItem]] = []
    current: list[TextItem] = []
    current_chars = 0
    for item in items:
        item_chars = len(item.text)
        if current and current_chars + item_chars > max_chars:
            batches.append(current)
            current = []
            current_chars = 0
        current.append(item)
        current_chars += item_chars
    if current:
        batches.append(current)
    return batches


def _extract_json_object(raw_text: str) -> dict[str, Any]:
    text = raw_text.strip()
    if text.startswith("```"):
        text = re.sub(r"^```(?:json)?\\n", "", text)
        text = re.sub(r"\\n```$", "", text)
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        match = re.search(r"\{.*\}", raw_text, flags=re.DOTALL)
        if not match:
            raise
        return json.loads(match.group(0))


def apply_never_translate_guard(text: str, protected: list[str]) -> str:
    for token in sorted(protected, key=len, reverse=True):
        if token in text:
            continue
        lowered = token.lower()
        match = re.search(re.escape(lowered), text.lower())
        if match:
            text = text[: match.start()] + token + text[match.end() :]
    return text


def enforce_glossary(text: str, glossary: dict[str, str]) -> str:
    fixed = text
    for source, target in glossary.items():
        pattern = re.compile(rf"\b{re.escape(source)}\b", re.IGNORECASE)
        fixed = pattern.sub(target, fixed)
    return fixed


def translate_batch_with_anthropic(
    *,
    batch: list[TextItem],
    glossary: dict[str, Any],
    target_language: str,
    prompt_template: str,
    model: str,
    api_key: str,
    timeout_s: int,
) -> dict[int, str]:
    payload_items = [
        {
            "id": item.id,
            "text": item.text,
            "kind": item.kind,
            "role": item.role,
            "location": item.location,
        }
        for item in batch
    ]

    user_prompt = (
        f"{prompt_template}\n\n"
        f"Target language: {target_language}\n"
        "Return JSON only in this schema: {\"translations\":[{\"id\":int,\"text\":str}]}\n"
        f"Never-translate list: {json.dumps(glossary.get('never_translate', []), ensure_ascii=False)}\n"
        f"Glossary map: {json.dumps(glossary.get('glossary', {}), ensure_ascii=False)}\n"
        f"Items: {json.dumps(payload_items, ensure_ascii=False)}"
    )

    body = {
        "model": model,
        "max_tokens": 4000,
        "temperature": 0,
        "system": "You are a deterministic translation engine. Return only strict JSON.",
        "messages": [{"role": "user", "content": user_prompt}],
    }

    req = urllib.request.Request(
        "https://api.anthropic.com/v1/messages",
        data=json.dumps(body).encode("utf-8"),
        headers={
            "Content-Type": "application/json",
            "x-api-key": api_key,
            "anthropic-version": "2023-06-01",
        },
        method="POST",
    )

    with urllib.request.urlopen(req, timeout=timeout_s) as response:
        response_obj = json.loads(response.read().decode("utf-8"))

    content = response_obj["content"][0]["text"]
    parsed = _extract_json_object(content)
    out: dict[int, str] = {}
    for row in parsed.get("translations", []):
        if isinstance(row, dict) and isinstance(row.get("id"), int) and isinstance(row.get("text"), str):
            out[row["id"]] = row["text"]
    return out


def _restore_whitespace(original: str, translated: str) -> str:
    leading = original[: len(original) - len(original.lstrip())]
    trailing = original[len(original.rstrip()) :]
    return leading + translated.strip() + trailing


def split_translation_to_runs(translated_text: str, run_count: int) -> list[str]:
    parts = translated_text.split(RUN_SEPARATOR)
    if len(parts) == run_count:
        return parts
    if run_count == 1:
        return [translated_text]
    return [translated_text] + ["" for _ in range(run_count - 1)]


def apply_translations(prs: Presentation, translations: dict[int, str]) -> None:
    idx = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    originals = _nonempty_run_texts(paragraph)
                    if not originals:
                        continue
                    translated = translations.get(idx)
                    if translated is not None:
                        split_runs = split_translation_to_runs(translated, len(originals))
                        run_idx = 0
                        for run in paragraph.runs:
                            if run.text and run.text.strip():
                                run.text = _restore_whitespace(originals[run_idx], split_runs[run_idx])
                                run_idx += 1
                    idx += 1

            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            originals = _nonempty_run_texts(paragraph)
                            if not originals:
                                continue
                            translated = translations.get(idx)
                            if translated is not None:
                                split_runs = split_translation_to_runs(translated, len(originals))
                                run_idx = 0
                                for run in paragraph.runs:
                                    if run.text and run.text.strip():
                                        run.text = _restore_whitespace(originals[run_idx], split_runs[run_idx])
                                        run_idx += 1
                            idx += 1

        rels = sorted(slide.part.rels.values(), key=lambda rel: rel.rId)
        for rel in rels:
            if rel.reltype != _DIAGRAM_DATA_RELTYPE:
                continue
            root = etree.fromstring(rel.target_part.blob)
            changed = False
            for element in list(root.iter(f"{{{_A_NS}}}t")):
                if element.text and element.text.strip():
                    translated = translations.get(idx)
                    if translated is not None:
                        element.text = _restore_whitespace(element.text, translated)
                        changed = True
                    idx += 1
            if changed:
                rel.target_part._blob = etree.tostring(
                    root, xml_declaration=True, encoding="UTF-8", standalone=True
                )

        if slide.has_notes_slide:
            for paragraph in slide.notes_slide.notes_text_frame.paragraphs:
                originals = _nonempty_run_texts(paragraph)
                if not originals:
                    continue
                translated = translations.get(idx)
                if translated is not None:
                    split_runs = split_translation_to_runs(translated, len(originals))
                    run_idx = 0
                    for run in paragraph.runs:
                        if run.text and run.text.strip():
                            run.text = _restore_whitespace(originals[run_idx], split_runs[run_idx])
                            run_idx += 1
                idx += 1


def load_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def preflight(args: argparse.Namespace) -> None:
    missing = []
    if not args.input_pptx.exists():
        missing.append(f"input PPTX not found: {args.input_pptx}")
    if not args.glossary.exists():
        missing.append(f"glossary artifact not found: {args.glossary}")
    if not args.prompt_template.exists():
        missing.append(f"prompt template artifact not found: {args.prompt_template}")
    if not os.environ.get("ANTHROPIC_API_KEY"):
        missing.append("ANTHROPIC_API_KEY is not set")

    if missing:
        raise SystemExit("Preflight failed:\n- " + "\n- ".join(missing))


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Translate PPTX with automated batching pipeline.")
    parser.add_argument("--input-pptx", type=Path, required=True)
    parser.add_argument("--output-pptx", type=Path, required=True)
    parser.add_argument("--target-language", required=True, help="Target language name, e.g., Spanish")
    parser.add_argument("--lang-code", default="es", help="Output language code for metadata")
    parser.add_argument("--glossary", type=Path, default=Path("glossary_en_es.json"))
    parser.add_argument("--prompt-template", type=Path, default=Path("translation-prompt-template.md"))
    parser.add_argument("--report-json", type=Path, default=Path("translation_report.json"))
    parser.add_argument("--batch-char-limit", type=int, default=4000)
    parser.add_argument("--max-retries", type=int, default=2)
    parser.add_argument("--model", default="claude-3-5-sonnet-latest")
    parser.add_argument("--timeout-s", type=int, default=90)
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    preflight(args)

    glossary_data = load_json(args.glossary)
    prompt_template = args.prompt_template.read_text(encoding="utf-8")
    api_key = os.environ["ANTHROPIC_API_KEY"]

    prs = Presentation(str(args.input_pptx))
    items = collect_text_items(prs)
    batches = build_batches(items, max_chars=args.batch_char_limit)

    translations: dict[int, str] = {}
    batch_reports: list[dict[str, Any]] = []

    for batch_index, batch in enumerate(batches, start=1):
        attempts = 0
        pending_ids = {item.id for item in batch}
        errors: list[str] = []

        while pending_ids and attempts <= args.max_retries:
            attempts += 1
            try:
                pending_batch = [item for item in batch if item.id in pending_ids]
                batch_result = translate_batch_with_anthropic(
                    batch=pending_batch,
                    glossary=glossary_data,
                    target_language=args.target_language,
                    prompt_template=prompt_template,
                    model=args.model,
                    api_key=api_key,
                    timeout_s=args.timeout_s,
                )

                for item in pending_batch:
                    translated = batch_result.get(item.id)
                    if translated is None:
                        continue
                    translated = apply_never_translate_guard(
                        translated, glossary_data.get("never_translate", [])
                    )
                    translated = enforce_glossary(
                        translated, glossary_data.get("glossary", {})
                    )
                    translations[item.id] = translated

                pending_ids = {item.id for item in pending_batch if item.id not in translations}
                if pending_ids:
                    time.sleep(min(2**attempts, 5))
            except (urllib.error.URLError, urllib.error.HTTPError, json.JSONDecodeError, KeyError) as exc:
                errors.append(str(exc))
                time.sleep(min(2**attempts, 5))

        batch_reports.append(
            {
                "batch_index": batch_index,
                "item_ids": [item.id for item in batch],
                "attempts": attempts,
                "translated_count": len([item for item in batch if item.id in translations]),
                "missing_ids": sorted([item.id for item in batch if item.id not in translations]),
                "errors": errors,
            }
        )

    missing_final = sorted([item.id for item in items if item.id not in translations])
    if missing_final:
        raise SystemExit(f"Translation failed for IDs: {missing_final}")

    apply_translations(prs, translations)
    prs.save(str(args.output_pptx))

    report = {
        "schema_version": "1.0",
        "timestamp_utc": datetime.now(timezone.utc).isoformat(),
        "input_pptx": str(args.input_pptx),
        "output_pptx": str(args.output_pptx),
        "target_language": args.target_language,
        "lang_code": args.lang_code,
        "total_items": len(items),
        "total_batches": len(batches),
        "translated_items": len(translations),
        "batches": batch_reports,
    }
    args.report_json.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"Translated {len(items)} items into {args.output_pptx}")
    print(f"Report: {args.report_json}")


if __name__ == "__main__":
    main()
